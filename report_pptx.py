"""
ESL Safety Observation - PPTX Report Generator (format.pptx style)
White background | Logo top-right | 3 observations per slide
Each observation = portrait photo + dark caption strip below.

Embeds photos directly in PPTX. No hyperlinks.
"""

import io
import os
import base64
from datetime import datetime

import requests
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE


# ---------- palette (matches format.pptx) ----------
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
NEAR_BLACK  = RGBColor(0x0D, 0x0D, 0x0D)   # caption fill in format.pptx
DARK_BLUE   = RGBColor(0x1F, 0x4E, 0x79)   # vedanta blue accent
ACCENT_GREEN = RGBColor(0x6E, 0xB4, 0x3F)  # vedanta green accent
LIGHT_GRAY  = RGBColor(0xF2, 0xF4, 0xF7)
MID_GRAY    = RGBColor(0xBF, 0xBF, 0xBF)
TEXT_DARK   = RGBColor(0x1A, 0x1A, 0x1A)
TEXT_MUTED  = RGBColor(0x59, 0x59, 0x59)

STATUS_COLORS = {
    "open":         RGBColor(0xDC, 0x26, 0x26),
    "in progress":  RGBColor(0xD9, 0x77, 0x06),
    "closed":       RGBColor(0x16, 0xA3, 0x4A),
}

AREAS = [
    "RMHS", "SINTER", "COKE OVEN", "POWERPLANT",
    "BLAST FURNACE 2", "BLAST FURNACE 3", "PCI", "PCM",
    "SECONDARY OPERATIONS",
]

# 16:9 widescreen
W = Inches(13.33)
H = Inches(7.5)

LOGO_PATH = os.path.join(os.path.dirname(os.path.abspath(__file__)), "logo.png")


# =============================================================
# Low-level helpers
# =============================================================
def _logo_stream():
    try:
        return open(LOGO_PATH, "rb")
    except Exception:
        return None


def _add_rect(slide, left, top, width, height, color, line=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if line:
        shape.line.color.rgb = color
    else:
        shape.line.fill.background()
    return shape


def _add_text(slide, text, left, top, width, height,
              size, bold=False, color=TEXT_DARK,
              align=PP_ALIGN.LEFT, wrap=True,
              font="Calibri"):
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame
    tf.word_wrap = wrap
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)

    p = tf.paragraphs[0]
    p.alignment = align

    run = p.add_run()
    run.text = str(text or "")
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = font
    run.font.color.rgb = color
    return tx


def _white_bg(slide):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE


def _status_color(status):
    s = (status or "open").lower()
    for key, col in STATUS_COLORS.items():
        if key in s:
            return col
    return STATUS_COLORS["open"]


# =============================================================
# Photo handling (unchanged behavior, embeds directly)
# =============================================================
def _photo_url_from_obs(obs):
    for key in ["image_url", "photo_url", "photo", "photo_link",
                "image", "image_link", "telegram_photo_url"]:
        val = obs.get(key)
        if val and str(val).startswith("http"):
            return str(val)
    return ""


def _fetch_image(obs):
    """Return a clean JPEG stream for python-pptx (embedded, not hyperlinked)."""
    try:
        data_url = obs.get("image_b64") or obs.get("closure_b64") or ""
        if data_url:
            if "," in data_url:
                data_url = data_url.split(",", 1)[1]
            raw = base64.b64decode(data_url)
        else:
            photo_url = _photo_url_from_obs(obs)
            if not photo_url:
                return None
            response = requests.get(photo_url, timeout=20)
            response.raise_for_status()
            raw = response.content

        img = Image.open(io.BytesIO(raw))
        img = img.convert("RGB")
        img.thumbnail((1600, 1600), Image.LANCZOS)

        out = io.BytesIO()
        img.save(out, format="JPEG", quality=88, optimize=True)
        out.seek(0)
        return out
    except Exception as e:
        print(f"Photo embed failed: {e}")
        return None


def _add_logo(slide, left, top, width, height):
    logo = _logo_stream()
    if logo:
        try:
            slide.shapes.add_picture(logo, left, top, width, height)
        except Exception:
            pass
        finally:
            logo.close()


def _add_photo_fit(slide, obs, left, top, width, height):
    """Add a photo; if missing, draw a gray placeholder of the same box."""
    img_data = _fetch_image(obs)
    if img_data:
        try:
            slide.shapes.add_picture(img_data, left, top, width, height)
            return
        except Exception as e:
            print(f"PPT photo add failed: {e}")
            label = "Photo Error"
    else:
        label = "No Photo"

    _add_rect(slide, left, top, width, height, MID_GRAY)
    _add_text(
        slide, label,
        left, top + height / 2 - Inches(0.2),
        width, Inches(0.4),
        size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
    )


# =============================================================
# Title slide (format.pptx style: white BG, logo top-right,
# large dark-blue title centered, accent line, stats at bottom)
# =============================================================
def _add_title_slide(prs, title_line1, title_line2, obs_list):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _white_bg(slide)

    # Logo, top-right (mirrors format.pptx: 3.03" x 0.58" near x=10.30, y=0.0)
    _add_logo(slide, Inches(10.20), Inches(0.20), Inches(3.00), Inches(0.65))

    # Thin separator line under header band (blue, full width)
    _add_rect(slide, Inches(0.32), Inches(1.05), W - Inches(0.64),
              Inches(0.025), DARK_BLUE)

    # Big centered title
    _add_text(
        slide, title_line1,
        Inches(0.5), Inches(2.4), Inches(12.33), Inches(1.4),
        size=44, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER,
        font="Calibri",
    )

    # Subtitle / date
    _add_text(
        slide, title_line2,
        Inches(0.5), Inches(3.75), Inches(12.33), Inches(0.7),
        size=24, color=TEXT_MUTED, align=PP_ALIGN.CENTER,
        font="Calibri",
    )

    # Stats row
    total = len(obs_list)
    open_c = sum(
        1 for o in obs_list
        if "open" in o.get("status", "").lower()
        and "progress" not in o.get("status", "").lower()
    )
    prog_c = sum(1 for o in obs_list if "progress" in o.get("status", "").lower())
    closed_c = sum(1 for o in obs_list if "closed" in o.get("status", "").lower())

    stats = [
        ("Total",       str(total),   DARK_BLUE),
        ("Open",        str(open_c),  STATUS_COLORS["open"]),
        ("In Progress", str(prog_c),  STATUS_COLORS["in progress"]),
        ("Closed",      str(closed_c), STATUS_COLORS["closed"]),
    ]

    box_w = Inches(2.6)
    box_h = Inches(1.4)
    gap = Inches(0.2)
    total_w = box_w * 4 + gap * 3
    start_x = (W - total_w) / 2
    y_box = Inches(4.95)

    for i, (label, val, col) in enumerate(stats):
        bx = start_x + i * (box_w + gap)

        # Top color strip (8% of card height) - matches format.pptx accent style
        _add_rect(slide, bx, y_box, box_w, Inches(0.10), col)
        # Card body
        _add_rect(slide, bx, y_box + Inches(0.10), box_w, box_h - Inches(0.10),
                  LIGHT_GRAY)

        _add_text(
            slide, val,
            bx, y_box + Inches(0.18),
            box_w, Inches(0.75),
            size=40, bold=True, color=col, align=PP_ALIGN.CENTER,
        )
        _add_text(
            slide, label.upper(),
            bx, y_box + Inches(0.92),
            box_w, Inches(0.4),
            size=14, bold=True, color=TEXT_DARK, align=PP_ALIGN.CENTER,
        )

    # Bottom green accent line (matches format.pptx green band)
    _add_rect(slide, Inches(0.32), Inches(6.95), W - Inches(0.64),
              Inches(0.03), ACCENT_GREEN)

    # Footer
    _add_text(
        slide,
        f"Generated: {datetime.now().strftime('%d/%m/%Y %H:%M')}   |   ESL Steel Limited",
        Inches(0.3), Inches(7.05), W - Inches(0.6), Inches(0.35),
        size=10, color=TEXT_MUTED, align=PP_ALIGN.CENTER,
    )


# =============================================================
# Content slide (format.pptx style: white BG, logo top-right,
# THREE portrait photos in a row, dark caption strip below each)
# =============================================================
def _add_obs_panel(slide, obs, left, top, photo_w, photo_h,
                   cap_left, cap_w, cap_top, cap_h):
    """One photo + caption block."""
    # Photo
    _add_photo_fit(slide, obs, left, top, photo_w, photo_h)

    # Tiny status pip on top-left corner of photo
    status = obs.get("status", "Open")
    sc = _status_color(status)
    pip_w = Inches(0.95)
    pip_h = Inches(0.28)
    _add_rect(slide, left + Inches(0.12), top + Inches(0.12),
              pip_w, pip_h, sc)
    _add_text(
        slide, status.upper(),
        left + Inches(0.12), top + Inches(0.12),
        pip_w, pip_h,
        size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
    )

    # Caption strip (matches format.pptx: dark fill, white text, centered)
    _add_rect(slide, cap_left, cap_top, cap_w, cap_h, NEAR_BLACK)

    # Caption text
    tx = slide.shapes.add_textbox(
        cap_left, cap_top, cap_w, cap_h,
    )
    tf = tx.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(6)
    tf.margin_right = Pt(6)
    tf.margin_top = Pt(3)
    tf.margin_bottom = Pt(3)

    # Line 1: observation (main caption text — like format.pptx)
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run()
    r1.text = str(obs.get("observation", "-"))
    r1.font.size = Pt(13)
    r1.font.bold = True
    r1.font.color.rgb = WHITE
    r1.font.name = "Calibri"

    # Line 2: meta (ref / area / target) — small, muted
    ref = obs.get("ref_no", "")
    area = obs.get("area", "")
    target = obs.get("target_date", "")
    meta_parts = [p for p in [ref, area, f"Target: {target}" if target else ""]
                  if p]
    if meta_parts:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = "   ·   ".join(meta_parts)
        r2.font.size = Pt(9)
        r2.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
        r2.font.name = "Calibri"


def _add_content_slide(prs, obs_group, area_label, page_num, page_total):
    """3 observations per slide, format.pptx layout."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _white_bg(slide)

    # Logo, top-right (same as format.pptx: roughly 3" x 0.58" at x=10.30)
    _add_logo(slide, Inches(10.30), Inches(0.00), Inches(3.03), Inches(0.58))

    # Optional small area label, top-left (subtle, not a big colored band)
    _add_text(
        slide, area_label.upper(),
        Inches(0.32), Inches(0.18), Inches(7), Inches(0.4),
        size=12, bold=True, color=DARK_BLUE, align=PP_ALIGN.LEFT,
    )
    _add_text(
        slide, f"Page {page_num} / {page_total}",
        W - Inches(3.35), Inches(0.62), Inches(3), Inches(0.3),
        size=9, color=TEXT_MUTED, align=PP_ALIGN.RIGHT,
    )

    # === Photo + caption geometry (mirrors format.pptx) ===
    # Photos: 3.84" x 5.12", at y=0.94, x = 0.32, 4.75, 9.17
    # Captions: 4.33" x 0.58", at y=6.41, x = 0.08, 4.50, 8.93
    photo_w = Inches(3.84)
    photo_h = Inches(5.12)
    photo_y = Inches(0.94)
    photo_xs = [Inches(0.32), Inches(4.75), Inches(9.17)]

    cap_w = Inches(4.33)
    cap_h = Inches(0.62)
    cap_y = Inches(6.41)
    cap_xs = [Inches(0.08), Inches(4.50), Inches(8.93)]

    for i, obs in enumerate(obs_group):
        if obs is None:
            continue
        _add_obs_panel(
            slide, obs,
            left=photo_xs[i], top=photo_y,
            photo_w=photo_w, photo_h=photo_h,
            cap_left=cap_xs[i], cap_w=cap_w,
            cap_top=cap_y, cap_h=cap_h,
        )

    # Thin green footer accent (matches format.pptx)
    _add_rect(slide, Inches(0.32), Inches(7.20), W - Inches(0.64),
              Inches(0.02), ACCENT_GREEN)


# =============================================================
# Section divider slide (optional - inserted between areas)
# =============================================================
def _add_section_slide(prs, area_label, count):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _white_bg(slide)

    _add_logo(slide, Inches(10.30), Inches(0.20), Inches(3.00), Inches(0.58))

    # Big area name centered
    _add_text(
        slide, area_label.upper(),
        Inches(0.5), Inches(2.8), Inches(12.33), Inches(1.4),
        size=54, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER,
    )
    _add_text(
        slide, f"{count} observation{'s' if count != 1 else ''}",
        Inches(0.5), Inches(4.2), Inches(12.33), Inches(0.7),
        size=22, color=TEXT_MUTED, align=PP_ALIGN.CENTER,
    )

    _add_rect(slide, Inches(0.32), Inches(6.95), W - Inches(0.64),
              Inches(0.03), ACCENT_GREEN)


# =============================================================
# Public API
# =============================================================
def generate_pptx(obs_list, area="ALL", date_label="", include_section_slide=False):
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    area_display = area if area != "ALL" else "All Areas"
    subtitle = date_label or datetime.now().strftime("%d/%m/%Y")

    # Title slide
    _add_title_slide(
        prs,
        f"Safety Observations — {area_display}",
        subtitle,
        obs_list,
    )

    if include_section_slide and obs_list:
        _add_section_slide(prs, area_display, len(obs_list))

    # 3 observations per slide
    chunk = 3
    page_total = (len(obs_list) + chunk - 1) // chunk
    for idx, i in enumerate(range(0, len(obs_list), chunk), start=1):
        group = obs_list[i:i + chunk]
        # pad to length 3 with None so geometry stays stable
        while len(group) < chunk:
            group.append(None)
        _add_content_slide(prs, group, area_display, idx, page_total)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def generate_pptx_per_area(obs_list, date_label=""):
    result = {}
    for area in AREAS:
        filtered = [
            o for o in obs_list
            if area.lower() in o.get("area", "").lower()
        ]
        if filtered:
            result[area] = generate_pptx(
                filtered,
                area=area,
                date_label=date_label,
            )
    return result
